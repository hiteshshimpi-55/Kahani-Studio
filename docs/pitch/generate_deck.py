#!/usr/bin/env python3
"""Kahani pitch deck — product-first: built → stack → problem → how it works → agents."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUT = ROOT / "Kahani_ZeroToOne_Pitch.pptx"

BG = RGBColor(0x0F, 0x0F, 0x12)
SURFACE = RGBColor(0x18, 0x18, 0x1F)
SURFACE_2 = RGBColor(0x22, 0x22, 0x2C)
ACCENT = RGBColor(0xE6, 0x19, 0x4D)
MUTED = RGBColor(0xA1, 0xA1, 0xAA)
DIM = RGBColor(0x71, 0x71, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0x2E, 0x2E, 0x38)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MX = Inches(0.70)
CONTENT_TOP = Inches(1.82)
FOOTER_Y = Inches(7.08)
CONTENT_W = Inches(11.93)
GAP = Inches(0.20)
TOTAL_PAGES = 8
FONT = "Calibri"


def _font(run, size, *, bold=False, color=WHITE):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = rPr.makeelement(qn(f"a:{tag}"), {})
            rPr.append(el)
        el.set("typeface", FONT)


def _bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _rect(slide, l, t, w, h, *, fill=SURFACE, line=None, rounded=False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    sh = slide.shapes.add_shape(kind, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if rounded:
        sh.adjustments[0] = 0.05
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.0)
    return sh


def _write(shape, blocks, *, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, ml=0.14, mr=0.12, mt=0.10, mb=0.08):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = valign
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    for i, (text, size, bold, color, space_after) in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(space_after)
        p.line_spacing = 1.1
        run = p.add_run()
        run.text = text
        _font(run, size, bold=bold, color=color)


def _label(slide, l, t, w, h, blocks, *, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    _write(box, blocks, align=align, ml=0, mr=0, mt=0, mb=0)
    return box


def _card(slide, l, t, w, h, blocks, *, fill=SURFACE, line=LINE, accent=False, center=False):
    align = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    valign = MSO_ANCHOR.MIDDLE if center else MSO_ANCHOR.TOP
    if accent:
        strip = Inches(0.08)
        _rect(slide, l, t, strip, h, fill=ACCENT)
        card = _rect(slide, l + strip, t, w - strip, h, fill=fill, line=line, rounded=True)
        _write(card, blocks, align=align, valign=valign, ml=0.12)
    else:
        card = _rect(slide, l, t, w, h, fill=fill, line=line, rounded=True)
        _write(card, blocks, align=align, valign=valign)
    return card


def _header(slide, section: str, title: str, subtitle: str | None = None):
    _rect(slide, MX, Inches(0.36), Inches(0.36), Inches(0.06), fill=ACCENT)
    _label(slide, MX, Inches(0.50), CONTENT_W, Inches(0.22), [(section.upper(), 11, True, ACCENT, 0)])
    _label(slide, MX, Inches(0.76), CONTENT_W, Inches(0.38), [(title, 22, True, WHITE, 0)])
    if subtitle:
        _label(slide, MX, Inches(1.20), CONTENT_W, Inches(0.40), [(subtitle, 12, False, MUTED, 0)])


def _footer(slide, page: int):
    _label(
        slide,
        MX,
        FOOTER_Y,
        Inches(9.2),
        Inches(0.22),
        [("Kahani  ·  Zero to One × Pocket FM  ·  IIM Bangalore", 10, False, DIM, 0)],
    )
    _label(
        slide,
        Inches(11.15),
        FOOTER_Y,
        Inches(1.5),
        Inches(0.22),
        [(f"{page:02d} / {TOTAL_PAGES:02d}", 10, False, DIM, 0)],
        align=PP_ALIGN.RIGHT,
    )


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    return s


def bullets(title: str, items: list[str], *, title_size=13, item_size=11):
    blocks = [(title, title_size, True, WHITE, 7)]
    for i, item in enumerate(items):
        gap = 4 if i < len(items) - 1 else 0
        blocks.append((f"•  {item}", item_size, False, MUTED, gap))
    return blocks


# ── Slides ──────────────────────────────────────────────────────────────────


def slide_01_cover(prs):
    s = blank(prs)
    _rect(s, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill=ACCENT)
    _label(
        s,
        MX,
        Inches(0.70),
        CONTENT_W,
        Inches(0.26),
        [("ZERO TO ONE  ·  POCKET FM × OPENAI × LIGHTSPEED  ·  IIM BANGALORE", 11, True, ACCENT, 0)],
    )
    _label(s, MX, Inches(2.05), CONTENT_W, Inches(0.75), [("Kahani", 52, True, WHITE, 0)])
    _label(
        s,
        MX,
        Inches(2.90),
        Inches(11.2),
        Inches(0.45),
        [("Multi-agent production studio for serialized audio entertainment.", 18, False, MUTED, 0)],
    )
    _label(
        s,
        MX,
        Inches(3.45),
        Inches(11.0),
        Inches(0.40),
        [
            (
                "What we built · how the agents work · Databricks-backed context & cast retrieval.",
                13,
                False,
                DIM,
                0,
            )
        ],
    )
    chips = [
        ("Theme", "T-04 Creator Tools  ·  T-01 / T-02"),
        ("Core", "7 specialized agents in one pipeline"),
        ("Demo", "Listen → simulate → human publish"),
    ]
    chip_w = Inches(3.80)
    for i, (a, b) in enumerate(chips):
        _card(
            s,
            MX + i * (chip_w + GAP),
            Inches(5.05),
            chip_w,
            Inches(1.20),
            [(a.upper(), 10, True, ACCENT, 6), (b, 12, True, WHITE, 0)],
        )
    return s


def slide_02_built(prs):
    s = blank(prs)
    _header(
        s,
        "01  What we built",
        "An end-to-end studio — from story hook to publishable episode package",
        "Not a chat wrapper: human-gated production stages with agent orchestration.",
    )
    left = [
        "Project workspace with context attachments",
        "Agent chat: discover, pitch, generate, approve",
        "LangGraph run → structured script package",
        "Voice + SFX mix (ElevenLabs / Sarvam)",
        "Companion visuals + cover art",
        "Web timeline editor (stems, markers, listen)",
        "Audience simulation (audit + personas + patches)",
        "Episode assembly → S3 artifact package",
    ]
    right = [
        ("Entry", "Project + prompt + attachments"),
        ("Gates", "Script → Audio → Visuals → Cover → Assembly"),
        ("Languages", "Hindi + English (library voices)"),
        ("Runtime", "Docker Compose · API · Worker · Redis"),
    ]
    _card(
        s,
        MX,
        CONTENT_TOP,
        Inches(6.40),
        Inches(4.55),
        bullets("Shipped surfaces", left, title_size=14, item_size=12),
        accent=True,
    )
    y = CONTENT_TOP
    rh = Inches(1.00)
    for i, (t, d) in enumerate(right):
        _card(
            s,
            MX + Inches(6.40) + GAP,
            y + i * (rh + Inches(0.12)),
            Inches(5.30),
            rh,
            [(t.upper(), 10, True, ACCENT, 3), (d, 13, True, WHITE, 0)],
        )
    _footer(s, 2)
    return s


def slide_03_stack(prs):
    s = blank(prs)
    _header(
        s,
        "02  Tech stack",
        "Production stack — including Databricks for vector retrieval",
        "Same stack we run in Compose locally and wire for cloud deploy.",
    )
    cols = [
        (
            "Application",
            [
                "React 19 · TypeScript · Vite · Tailwind",
                "FastAPI · SQLAlchemy async · Alembic",
                "ARQ workers on Redis",
                "LangGraph + Postgres checkpointer",
            ],
        ),
        (
            "AI & media",
            [
                "LLM: OpenAI / Anthropic (Script Writer)",
                "TTS: ElevenLabs + Sarvam (Hindi)",
                "Images: OpenAI · Gemini (fallback)",
                "Research: Tavily web crawl",
            ],
        ),
        (
            "Data & infra",
            [
                "Postgres 16 · Redis 7",
                "Databricks AI Search (project RAG)",
                "Databricks Vector Search (cast catalog)",
                "S3 artifacts · Docker · AWS / Terraform",
            ],
        ),
    ]
    card_w = Inches(3.80)
    for i, (title, items) in enumerate(cols):
        _card(
            s,
            MX + i * (card_w + GAP),
            CONTENT_TOP,
            card_w,
            Inches(3.55),
            bullets(title, items, title_size=13, item_size=11),
            accent=True,
        )

    _card(
        s,
        MX,
        CONTENT_TOP + Inches(3.75),
        CONTENT_W,
        Inches(0.95),
        [
            ("Databricks in the loop", 12, True, ACCENT, 4),
            (
                "AI Search indexes story attachments for grounded scripting. Vector Search retrieves cast / voice assets and shot templates — so agents don’t invent casting or visuals from thin air.",
                12,
                False,
                MUTED,
                0,
            ),
        ],
        line=ACCENT,
    )
    _footer(s, 3)
    return s


def slide_04_problem(prs):
    s = blank(prs)
    _header(
        s,
        "03  Problem",
        "Serial audio wins on part-to-part retention. Production is still fragmented.",
        "Pocket FM–class storytelling needs speed, cultural specificity, and pre-publish signal.",
    )
    pains = [
        ("Idea → script is slow", "Regional authenticity is the moat. Generic LLM drafts sound thin."),
        ("Episodes don’t chain", "Weak cliffhangers → drop-off after part one."),
        ("Audio assembly is manual", "Voice, SFX, and visuals are synced by hand across tools."),
        ("No pre-publish signal", "Teams guess which cohorts will continue listening."),
        ("No closed loop", "Wins on one story don’t systematically improve the next."),
        ("Tool sprawl", "Briefing, writing, casting, mix, and QA live in five apps."),
    ]
    card_w, card_h = Inches(3.80), Inches(1.85)
    for i, (t, d) in enumerate(pains):
        col, row = i % 3, i // 3
        _card(
            s,
            MX + col * (card_w + GAP),
            CONTENT_TOP + row * (card_h + GAP),
            card_w,
            card_h,
            [(t, 13, True, WHITE, 6), (d, 11, False, MUTED, 0)],
            accent=True,
        )
    _footer(s, 4)
    return s


def slide_05_how(prs):
    s = blank(prs)
    _header(
        s,
        "04  How it works",
        "One orchestrated pipeline — agents specialize, humans gate quality cliffs",
        "Auto-run generation. Never auto-publish.",
    )

    stages = [
        ("1", "Discover"),
        ("2", "Research"),
        ("3", "Script"),
        ("4", "Narrate"),
        ("5", "Voice"),
        ("6", "Visuals"),
        ("7", "Edit"),
        ("8", "Simulate"),
    ]
    sw = Inches(1.35)
    for i, (n, name) in enumerate(stages):
        _card(
            s,
            MX + i * (sw + Inches(0.12)),
            CONTENT_TOP,
            sw,
            Inches(1.15),
            [(n, 12, True, ACCENT, 2), (name, 11, True, WHITE, 0)],
            center=True,
        )

    flow = [
        (
            "Ground",
            "Attachments → Databricks AI Search. Cast & templates → Vector Search. Optional Tavily crawl for fresh hooks.",
        ),
        (
            "Write",
            "Storytelling + Script Writer produce a structured multi-part package with cliffhangers and narration plan.",
        ),
        (
            "Produce",
            "Voice agent renders TTS + SFX. Director + Image agents build companion visuals and cover art.",
        ),
        (
            "Assure",
            "Timeline editor for listen/edit. Audience Sim proposes patches. Human approves each gate.",
        ),
    ]
    fw = Inches(2.80)
    y = CONTENT_TOP + Inches(1.40)
    for i, (t, d) in enumerate(flow):
        _card(
            s,
            MX + i * (fw + GAP),
            y,
            fw,
            Inches(2.95),
            [(t.upper(), 11, True, ACCENT, 8), (d, 12, False, WHITE, 0)],
        )
    _footer(s, 5)
    return s


def slide_06_agents(prs):
    s = blank(prs)
    _header(
        s,
        "05  Agents",
        "Seven specialized agents — one production studio",
        "Each owns a stage; the orchestrator sequences them with human approvals.",
    )
    agents = [
        ("Storytelling Agent", "Shapes the series brief, plot pitches, and creative direction inside chat."),
        ("Scripting & Search Agent", "Retrieves project context (Databricks) and researches hooks (Tavily + extraction)."),
        ("Narrative Agent", "Locks narration mode — narration-led, dialogue inserts, multi-narrator, framed."),
        ("Script Writer Agent", "Writes the structured multi-part script package and screenplay."),
        ("Director Agent", "Plans visual shots / lookbook — performance-aware scene direction."),
        ("Voice Agent", "Casts library voices and renders TTS + SFX beds (ElevenLabs / Sarvam)."),
        ("Image Agent", "Generates companion stills, cover art, and visual track assets."),
    ]
    # 4 on top row, 3 on bottom — balanced
    card_w, card_h = Inches(2.80), Inches(1.95)
    for i, (name, desc) in enumerate(agents[:4]):
        _card(
            s,
            MX + i * (card_w + GAP),
            CONTENT_TOP,
            card_w,
            card_h,
            [(name, 12, True, WHITE, 6), (desc, 11, False, MUTED, 0)],
            accent=True,
        )
    # bottom 3 centered-ish across width
    bottom_w = Inches(3.80)
    start_x = MX + Inches(0.15)
    for i, (name, desc) in enumerate(agents[4:]):
        _card(
            s,
            start_x + i * (bottom_w + GAP),
            CONTENT_TOP + card_h + GAP,
            bottom_w,
            card_h,
            [(name, 12, True, WHITE, 6), (desc, 11, False, MUTED, 0)],
            accent=True,
        )
    _footer(s, 6)
    return s


def slide_07_flow_detail(prs):
    s = blank(prs)
    _header(
        s,
        "06  Agent handoff",
        "Who does what in sequence",
        "Search grounds truth → writers lock story → voice & image produce → humans decide.",
    )
    rows = [
        ("Storytelling", "Chat director", "Plot pitches, series intent, route generate / rewrite"),
        ("Scripting & Search", "RAG + crawl", "Databricks chunks · Tavily research · SOURCE assembly"),
        ("Narrative + Script Writer", "Script package", "Narration config + multi-part screenplay + cliffs"),
        ("Voice", "Audio stems", "Cast match (Vector Search) · TTS · SFX mix"),
        ("Director + Image", "Visual track", "Shot plan · lookbook · stills · cover art"),
        ("Human + Audience Sim", "Quality gate", "Approve stages · listen in editor · apply patches"),
    ]
    row_h = Inches(0.68)
    for i, (agent, out, detail) in enumerate(rows):
        y = CONTENT_TOP + i * (row_h + Inches(0.08))
        _card(
            s,
            MX,
            y,
            Inches(3.10),
            row_h,
            [(agent, 12, True, WHITE, 0)],
            fill=SURFACE_2,
        )
        _card(
            s,
            MX + Inches(3.20),
            y,
            Inches(2.40),
            row_h,
            [(out, 11, True, ACCENT, 0)],
        )
        _card(
            s,
            MX + Inches(5.70),
            y,
            Inches(6.20),
            row_h,
            [(detail, 11, False, MUTED, 0)],
        )
    _footer(s, 7)
    return s


def slide_08_close(prs):
    s = blank(prs)
    _rect(s, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill=ACCENT)
    _label(s, MX, Inches(1.35), CONTENT_W, Inches(0.26), [("THANK YOU", 12, True, ACCENT, 0)])
    _label(s, MX, Inches(1.85), CONTENT_W, Inches(0.70), [("Kahani", 46, True, WHITE, 0)])
    _label(
        s,
        MX,
        Inches(2.65),
        Inches(11.0),
        Inches(0.55),
        [
            (
                "Seven agents. One timeline. Human publish. Built for Pocket FM–scale serial production.",
                15,
                False,
                MUTED,
                0,
            )
        ],
    )
    asks = [
        ("Ask", "Pilot with a Pocket FM–style content pod"),
        ("Proof", "Live demo: generate → listen → simulate"),
        ("Contact", "Add team names and emails here"),
    ]
    card_w = Inches(3.80)
    for i, (a, b) in enumerate(asks):
        _card(
            s,
            MX + i * (card_w + GAP),
            Inches(3.70),
            card_w,
            Inches(1.35),
            [(a.upper(), 11, True, ACCENT, 6), (b, 13, True, WHITE, 0)],
        )
    _label(
        s,
        MX,
        Inches(5.50),
        CONTENT_W,
        Inches(0.70),
        [
            ("Stack highlight", 11, True, WHITE, 4),
            (
                "LangGraph · FastAPI · React · ElevenLabs/Sarvam · Gemini/OpenAI · Tavily · Databricks AI Search & Vector Search · S3 · Docker",
                12,
                False,
                DIM,
                0,
            ),
        ],
    )
    _footer(s, 8)
    return s


def main():
    ASSETS.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_cover(prs)
    slide_02_built(prs)
    slide_03_stack(prs)
    slide_04_problem(prs)
    slide_05_how(prs)
    slide_06_agents(prs)
    slide_07_flow_detail(prs)
    slide_08_close(prs)

    try:
        prs.save(OUT)
        print(f"Wrote {OUT} ({len(prs.slides)} slides)")
    except PermissionError:
        alt = ROOT / "Kahani_ZeroToOne_Pitch_v2.pptx"
        prs.save(alt)
        print(f"Original locked; wrote {alt} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
